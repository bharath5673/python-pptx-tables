from pptx import Presentation
from pptx.util import Inches, Pt
from pptx_tables import PptxTable

prs=Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


data = [
        [['Aa', 'Bb','Cc'],
        [3,     4,      5],
        [6,     7,      8]],

        [['Dd', 'Ee','Ff'],
        [0,     1,      2],
        [6,     7,      8]]
        ]



for n,lst in enumerate(data):

    lyt=prs.slide_layouts[6] # choosing a slide layout
    slide=prs.slides.add_slide(lyt) # adding a slide
    # title=slide.shapes.title 
    title_name = f"Page : "+str(n)
    # title.text=title_name
    # subtitle=slide.placeholders[1]


    tbl = PptxTable(lst,prs)
    tbl.set_table_location(left=Inches(1.5), top=Inches(2), width=Inches(5))
    # tbl.set_formatting(font_size=Pt(7), row_height=Inches(.3),alignment=PP_PARAGRAPH_ALIGNMENT.LEFT)
    tbl.set_formatting(font_size=Pt(12), row_height=Inches(.85))        
    tbl.create_table(slide_index=n,
                      columns_widths_weight=[2, 2, 2],
                      transpose=True,
                      )


tbl.save_pptx("slide_table.pptx")